from fastapi import FastAPI, File, UploadFile, HTTPException
import fitz  # PyMuPDF for PDF
from docx import Document  # for DOCX
from pptx import Presentation  # for PPTX
from io import BytesIO

app = FastAPI()

@app.post("/read-file/")
async def read_file(file: UploadFile = File(...)):
    content_type = file.content_type
    file_bytes = await file.read()

    # --- PDF ---
    if content_type == "application/pdf":
        text = extract_text_from_pdf(file_bytes)

    # --- DOCX ---
    elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = extract_text_from_docx(file_bytes)

    # --- PPTX ---
    elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = extract_text_from_pptx(file_bytes)

    else:
        raise HTTPException(status_code=400, detail="Only PDF, DOCX, or PPTX files are supported")

    return {"filename": file.filename, "content": text}



def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    text = ""
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text.strip()


def extract_text_from_docx(docx_bytes: bytes) -> str:
    text = ""
    doc = Document(BytesIO(docx_bytes))
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text.strip()


def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    text = ""
    prs = Presentation(BytesIO(pptx_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

from fastapi.middleware.cors import CORSMiddleware

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
